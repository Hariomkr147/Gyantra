"""PPTX parsing with python-pptx — slide titles become headings, body text follows."""

from __future__ import annotations

import logging

from app.models.enums import BlockType, ParseRoute
from app.models.schemas import DocumentIntelligenceResult, TextBlock

logger = logging.getLogger("gyantra.parser.pptx")


def parse_pptx(file_path: str, route: ParseRoute = ParseRoute.LIGHTWEIGHT_TEXT) -> DocumentIntelligenceResult:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    blocks: list[TextBlock] = []
    plain_parts: list[str] = []
    table_count = figure_count = 0

    for idx, slide in enumerate(prs.slides, start=1):
        # Slide title becomes an h2-level heading; deck itself is the h1.
        title_text = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
        except AttributeError:
            pass

        if title_text:
            blocks.append(
                TextBlock(
                    block_type=BlockType.HEADING,
                    content=title_text,
                    level=2,
                    page=idx,
                )
            )
            plain_parts.append(title_text)
        else:
            blocks.append(
                TextBlock(
                    block_type=BlockType.HEADING,
                    content=f"Slide {idx}",
                    level=2,
                    page=idx,
                )
            )

        for shape in slide.shapes:
            # Skip the title, already captured
            if title_text and shape == slide.shapes.title:
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                figure_count += 1
                blocks.append(
                    TextBlock(
                        block_type=BlockType.FIGURE,
                        content=f"[Image on slide {idx}]",
                        page=idx,
                    )
                )
                continue

            if shape.has_table:
                rows = [
                    [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                md = _rows_to_markdown(rows)
                blocks.append(TextBlock(block_type=BlockType.TABLE, content=md, page=idx))
                plain_parts.append(md)
                table_count += 1
                continue

            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                btype = BlockType.LIST if para.level > 0 else BlockType.PARAGRAPH
                blocks.append(TextBlock(block_type=btype, content=text, page=idx))
                plain_parts.append(text)

        # Speaker notes are often the richest teaching content in a deck.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    blocks.append(
                        TextBlock(
                            block_type=BlockType.PARAGRAPH,
                            content=f"[Speaker notes] {notes}",
                            page=idx,
                        )
                    )
                    plain_parts.append(notes)
        except (AttributeError, ValueError):
            pass

    return DocumentIntelligenceResult(
        blocks=blocks,
        plain_text="\n".join(plain_parts),
        page_count=len(prs.slides),
        file_type="pptx",
        parse_route=route,
        table_count=table_count,
        figure_count=figure_count,
        equation_count=0,
        metadata={"slide_count": len(prs.slides), "block_count": len(blocks)},
    )


def _rows_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    padded = [r + [""] * (width - len(r)) for r in rows]
    header, *body = padded
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)
